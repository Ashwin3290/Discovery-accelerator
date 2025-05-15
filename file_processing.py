import os
import io
import logging
import psutil
from typing import List, Dict, Any
from pathlib import Path
import fitz  # PyMuPDF for PDF processing
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import chromadb
import numpy as np
import tempfile
import gc
import requests
import base64
# from sentence_transformers import SentenceTransformer
import chromadb
from sklearn.cluster import KMeans
import numpy as np
import google.generativeai as genai
import os
from tqdm import tqdm
import docx  # Python-docx for DOCX processing
import zipfile
import xml.etree.ElementTree as ET
import time
import sqlite3
import pandas as pd
import os
import tempfile
from docx2pdf import convert as docx_convert
from pptxtopdf import convert
from comtypes import client
import google.generativeai as genai
from PIL import Image
import io
import re
import json
import gc

# Set up Gemini API
os.environ["GOOGLE_API_KEY"] = "AIzaSyAo-zfJCvUMxa91I5oC6r7AWSaWC6cn0cw"
genai.configure(api_key=os.environ["GOOGLE_API_KEY"])

# text_embedding_model = SentenceTransformer('sentence-transformers/all-MiniLM-L6-v2')
def get_memory_usage():
    """Get current memory usage of the process"""
    process = psutil.Process(os.getpid())
    return process.memory_info().rss / 1024 / 1024  # in MB

class ProjectDataPipeline:
    def __init__(self, base_dir: str, inference_api_url: str, chroma_path: str = None, gemini_api_key: str = None):
        # Store the inference API URL
        self.inference_api_url = inference_api_url.rstrip('/')
        self.gemini_api_key = gemini_api_key
        # Use temporary directory if no chroma_path provided
        if chroma_path is None:
            chroma_path = tempfile.mkdtemp()
        
        # Setup logging
        logging.basicConfig(level=logging.INFO)
        self.logger = logging.getLogger(__name__)
        self.logger.info(f"Initial RAM usage: {get_memory_usage():.2f} MB")
        
        # Initialize ChromaDB client
        os.makedirs(chroma_path, exist_ok=True)
        self.chroma_client = chromadb.PersistentClient(
            path=chroma_path,
            settings=chromadb.Settings(
                allow_reset=True,
                is_persistent=True
            )
        )
        
        self.base_dir = base_dir
        self.logger.info(f"Initialized pipeline with ChromaDB at: {chroma_path}")
        self.logger.info(f"Using inference API at: {inference_api_url}")
    
    def clear_memory(self):
        """Clear system memory"""
        gc.collect()
        memory_mb = get_memory_usage()
        self.logger.info(f"Current RAM usage: {memory_mb:.2f} MB")
    
    def preprocess_image(self, image: Image, max_size: int = 384) -> Image:
        """Resize image to reduce memory usage while maintaining aspect ratio"""
        try:
            if image.mode != 'RGB':
                image = image.convert('RGB')
            
            ratio = max_size / max(image.size)
            if ratio < 1:
                new_size = tuple(int(dim * ratio) for dim in image.size)
                image = image.resize(new_size, Image.Resampling.LANCZOS)
            
            return image
        except Exception as e:
            self.logger.error(f"Error preprocessing image: {str(e)}")
            return image
    
    def create_project_collection(self, project_name: str):
        """Create a new collection for a project if it doesn't exist"""
        try:
            collection = self.chroma_client.create_collection(name=project_name)
            self.logger.info(f"Created new collection for project: {project_name}")
            return collection
        except ValueError:
            collection = self.chroma_client.get_collection(name=project_name)
            self.logger.info(f"Collection already exists for project: {project_name}")
            return collection
    
    def process_image(self, image: Image) -> np.ndarray:
        """Process image using the inference API"""
        try:
            # Preprocess and resize image
            image = self.preprocess_image(image)
            
            # Convert image to base64
            buffered = io.BytesIO()
            image.save(buffered, format="PNG")
            img_str = base64.b64encode(buffered.getvalue()).decode()
            
            # Call the inference API
            response = requests.post(
                f"{self.inference_api_url}/process_image",
                json={'image': img_str}
            )
            
            if response.status_code != 200:
                raise Exception(f"API Error: {response.text}")
            
            result = response.json()
            return np.array(result['embedding'])
        
        except Exception as e:
            self.logger.error(f"Error in process_image: {str(e)}")
            return np.zeros((1, 384), dtype=np.float32)
    
    def get_text_embedding(self, text: str) -> np.ndarray:
        """Get text embedding from the inference API"""
        try:
            print(f"Getting text embedding for text of length {len(text)}")
            
            # Try first with the API
            try:
                print(f"Calling inference API at {self.inference_api_url}/embed_text")
                response = requests.post(
                    f"{self.inference_api_url}/embed_text",
                    json={'text': text}
                )
                
                if response.status_code == 200:
                    result = response.json()
                    print("Successfully got embedding from API")
                    return np.array(result['embedding'])
                else:
                    print(f"API Error ({response.status_code}): {response.text}")
                    print("Falling back to local model")
            except Exception as api_e:
                print(f"Error calling API: {str(api_e)}")
                print("Falling back to local model")
            
            # Fallback to local model
            # print("Using local model for text embedding")
            # embedding = text_embedding_model.encode(
            #     text,
            #     convert_to_numpy=True,
            #     normalize_embeddings=True
            # )
            # print(f"Successfully generated embedding with shape {embedding.shape}")
            # return embedding
        except Exception as e:
            self.logger.error(f"Error getting text embedding: {str(e)}")
            return np.zeros((384,), dtype=np.float32)

    # def decode_text_embedding(self,embedding: np.ndarray) -> str:
    #     """Decode text embedding back to text using the text_embedding_model"""
    #     try:
    #         # Decode the embedding to text
    #         decoded_text = text_embedding_model.decode(embedding)
    #         return decoded_text
    #     except Exception as e:
    #         self.logger.error(f"Error decoding text embedding: {str(e)}")
    #         return ""
        
    def create_embeddings(self, documents: List[Dict[str, Any]]) -> Dict[str, List]:
        """Create embeddings for text and images with memory management"""
        print(f"Creating embeddings for {len(documents)} documents")
        embeddings = []
        metadatas = []
        ids = []
        
        batch_size = 10
        for idx in range(0, len(documents), batch_size):
            batch = documents[idx:idx + batch_size]
            print(f"Processing batch {idx//batch_size + 1}/{(len(documents)-1)//batch_size + 1}")
            
            for doc_idx, doc in enumerate(batch):
                try:
                    if doc['type'] == 'text':
                        print(f"Creating text embedding for document {idx + doc_idx} from source {doc['source']}")
                        text_content = doc['content']
                        if isinstance(text_content, str) and len(text_content) > 0:
                            embedding = self.get_text_embedding(text_content).tolist()
                        else:
                            print(f"Warning: Empty or invalid text content in document {idx + doc_idx}")
                            continue
                    else:  # image
                        print(f"Processing image embedding for document {idx + doc_idx} from source {doc['source']}")
                        embedding = doc['content'].flatten().tolist()
                    
                    if embedding and len(embedding) > 0:
                        embeddings.append(embedding)
                        metadatas.append({
                            'source': doc['source'],
                            'type': doc['type']
                        })
                        ids.append(f"doc_{idx + doc_idx}")
                        print(f"Successfully created embedding for document {idx + doc_idx}")
                    else:
                        print(f"Warning: Skipping document {idx + doc_idx} due to empty embedding")
                        self.logger.warning(f"Skipping document {idx + doc_idx} due to empty embedding")
                except Exception as e:
                    print(f"Error creating embedding for document {idx + doc_idx}: {str(e)}")
                    self.logger.error(f"Error creating embedding for document {idx + doc_idx}: {str(e)}")
                    continue
            
            self.clear_memory()
        
        print(f"Successfully created {len(embeddings)} embeddings")
        
        result = {
            'embeddings': embeddings,
            'metadatas': metadatas,
            'ids': ids
        }
        
        # Verify result has data
        if not embeddings:
            print("WARNING: No embeddings were created!")
        
        return result
    

    def query_project(self, project_name: str, query_text: str, n_results: int = 5) -> List[Dict]:
        """Query a specific project's collection"""
        try:
            collection = self.chroma_client.get_collection(name=project_name)
            
            # Get query embedding from the API
            query_embedding = self.get_text_embedding(query_text).tolist()
            
            # Rest of the method remains the same as in your original code
            results = collection.query(
                query_embeddings=[query_embedding],
                n_results=n_results,
                include=["metadatas", "distances", "documents", "embeddings"]
            )
            
            # Format results (same as original code)
            formatted_results = []
            for idx in range(len(results['ids'][0])):
                result = {
                    'source': results['metadatas'][0][idx]['source'],
                    'type': results['metadatas'][0][idx]['type'],
                    'distance': results['distances'][0][idx],
                    'embedding': results['embeddings'][0][idx],
                    "Text": self.decode_text_embedding(results['embeddings'][0][idx])
                    
                }
                
                if results['metadatas'][0][idx]['type'] == 'text':
                    if idx < len(results['documents'][0]) and results['documents'][0][idx]:
                        result['content'] = results['documents'][0][idx]
                    else:
                        result['content'] = self.get_document_content(result['source'])
                else:
                    try:
                        if '#page' in result['source']:
                            pdf_path, page_info = result['source'].split('#page')
                            page_num, img_num = page_info.split('_img')
                            page_num, img_num = int(page_num) - 1, int(img_num)
                            
                            with fitz.open(pdf_path) as doc:
                                page = doc[page_num]
                                images = page.get_images()
                                if img_num < len(images):
                                    xref = images[img_num][0]
                                    base_image = doc.extract_image(xref)
                                    if base_image:
                                        image_data = base_image["image"]
                                        result['image'] = Image.open(io.BytesIO(image_data))
                        
                        elif '#slide' in result['source']:
                            # PPT image handling
                            ppt_path, slide_info = result['source'].split('#slide')
                            slide_num, img_num = slide_info.split('_img')
                            slide_num, img_num = int(slide_num) - 1, int(img_num)
                            
                            presentation = Presentation(ppt_path)
                            slide = presentation.slides[slide_num]
                            image_shapes = [shape for shape in slide.shapes
                                          if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
                            if img_num < len(image_shapes):
                                image_data = image_shapes[img_num].image.blob
                                result['image'] = Image.open(io.BytesIO(image_data))
                        
                        else:  # Regular image file
                            result['image'] = Image.open(result['source'])
                    
                    except Exception as e:
                        self.logger.error(f"Error loading image from {result['source']}: {str(e)}")
                        result['image'] = None
                
                formatted_results.append(result)
            
            return formatted_results
        
        except Exception as e:
            self.logger.error(f"Error querying project {project_name}: {str(e)}")
            return []

    # def extract_text_from_pdf(self, pdf_path: str) -> List[Dict[str, Any]]:
    #     """Extract text and images from PDF with memory management"""
    #     documents = []
        
    #     try:
    #         doc = fitz.open(pdf_path)
    #         total_pages = len(doc)
            
    #         for page_num in range(total_pages):
    #             self.clear_memory()
                
    #             page = doc[page_num]
                
    #             # Extract text
    #             text = page.get_text()
    #             if text.strip():
    #                 documents.append({
    #                     'content': text,
    #                     'type': 'text',
    #                     'source': f"{pdf_path}#page{page_num + 1}"
    #                 })
                
    #             # Process images for current page
    #             page_images = []
    #             page_sources = []
                
    #             images = page.get_images()
    #             for img_index, img in enumerate(images):
    #                 try:
    #                     xref = img[0]
    #                     base_image = doc.extract_image(xref)
    #                     if base_image:
    #                         image_data = base_image["image"]
    #                         image = Image.open(io.BytesIO(image_data))
                            
    #                         if image.size[0] >= 32 and image.size[1] >= 32:
    #                             page_images.append(image)
    #                             page_sources.append(f"{pdf_path}#page{page_num + 1}_img{img_index}")
    #                 except Exception as e:
    #                     self.logger.error(f"Error extracting image {img_index} from page {page_num + 1}: {str(e)}")
    #                     continue
                
    #             # Process images for current page
    #             if page_images:
    #                 image_documents = self.process_images_batch(page_images, page_sources)
    #                 documents.extend(image_documents)
                    
    #                 for img in page_images:
    #                     del img
    #                 del page_images
    #                 self.clear_memory()
                
    #             del page
    #             self.clear_memory()
            
    #         doc.close()
        
    #     except Exception as e:
    #         self.logger.error(f"Error processing PDF {pdf_path}: {str(e)}")
    #     finally:
    #         self.clear_memory()
        
    #     return documents

    # def extract_from_ppt(self, ppt_path: str) -> List[Dict[str, Any]]:
    #     """Extract text and images from PowerPoint with memory management"""
    #     documents = []
    #     try:
    #         presentation = Presentation(ppt_path)
            
    #         for slide_num, slide in enumerate(presentation.slides):
    #             self.clear_memory()
                
    #             # Extract text
    #             text = ""
    #             for shape in slide.shapes:
    #                 if hasattr(shape, "text"):
    #                     text += shape.text + "\n"
                
    #             if text.strip():
    #                 documents.append({
    #                     'content': text,
    #                     'type': 'text',
    #                     'source': f"{ppt_path}#slide{slide_num + 1}"
    #                 })
                
    #             # Process images for current slide
    #             slide_images = []
    #             slide_sources = []
                
    #             for shape_num, shape in enumerate(slide.shapes):
    #                 try:
    #                     if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
    #                         image = Image.open(io.BytesIO(shape.image.blob))
                            
    #                         if image.size[0] >= 32 and image.size[1] >= 32:
    #                             slide_images.append(image)
    #                             slide_sources.append(f"{ppt_path}#slide{slide_num + 1}_img{shape_num}")
                    
    #                 except Exception as e:
    #                     self.logger.error(f"Error extracting image {shape_num} from slide {slide_num + 1}: {str(e)}")
    #                     continue
                
    #             # Process slide images
    #             if slide_images:
    #                 image_documents = self.process_images_batch(slide_images, slide_sources)
    #                 documents.extend(image_documents)
                    
    #                 for img in slide_images:
    #                     del img
    #                 del slide_images
    #                 self.clear_memory()
                
    #             self.clear_memory()
        
    #     except Exception as e:
    #         self.logger.error(f"Error processing PPT {ppt_path}: {str(e)}")
    #     finally:
    #         self.clear_memory()
        
    #     return documents

    def process_images_batch(self, images: List[Image.Image], sources: List[str]) -> List[Dict]:
        """Process a batch of images with memory management"""
        documents = []
        batch_size = 1  # Process one at a time to manage memory
        
        for i in range(0, len(images), batch_size):
            current_batch = []
            current_sources = []
            
            for j in range(i, min(i + batch_size, len(images))):
                try:
                    img_copy = images[j].copy()
                    current_batch.append(img_copy)
                    current_sources.append(sources[j])
                except Exception as e:
                    self.logger.error(f"Error copying image {sources[j]}: {str(e)}")
                    continue
            
            for img, src in zip(current_batch, current_sources):
                try:
                    embedding = self.process_image(img)
                    documents.append({
                        'content': embedding,
                        'type': 'image',
                        'source': src
                    })
                except Exception as e:
                    self.logger.error(f"Error processing image {src}: {str(e)}")
                finally:
                    del img
                    self.clear_memory()
            
            del current_batch
            self.clear_memory()
        
        return documents

    # def extract_from_docx(self, docx_path: str) -> List[Dict[str, Any]]:
    #     """Extract text and images from DOCX with memory management"""
    #     documents = []
        
    #     try:
    #         doc = docx.Document(docx_path)
            
    #         # Extract text from paragraphs
    #         all_paragraphs = []
    #         for para in doc.paragraphs:
    #             if para.text.strip():
    #                 all_paragraphs.append(para.text)
            
    #         # Combine paragraphs into chunks for better context
    #         chunk_size = 3
    #         for i in range(0, len(all_paragraphs), chunk_size):
    #             chunk = all_paragraphs[i:i + chunk_size]
    #             text = "\n".join(chunk)
    #             if text.strip():
    #                 documents.append({
    #                     'content': text,
    #                     'type': 'text',
    #                     'source': f"{docx_path}#para{i}"
    #                 })
            
    #         # Extract tables
    #         for table_idx, table in enumerate(doc.tables):
    #             table_text = []
    #             for row in table.rows:
    #                 row_text = []
    #                 for cell in row.cells:
    #                     row_text.append(cell.text)
    #                 table_text.append(" | ".join(row_text))
                
    #             table_content = "\n".join(table_text)
    #             if table_content.strip():
    #                 documents.append({
    #                     'content': table_content,
    #                     'type': 'text',
    #                     'source': f"{docx_path}#table{table_idx}"
    #                 })
            
    #         # Extract images
    #         with zipfile.ZipFile(docx_path) as zf:
    #             image_parts = []
    #             for f in zf.namelist():
    #                 if f.startswith('word/media/'):
    #                     image_parts.append(f)
                
    #             for img_idx, img_part in enumerate(image_parts):
    #                 try:
    #                     image_content = zf.read(img_part)
    #                     image = Image.open(io.BytesIO(image_content))
                        
    #                     if image.size[0] >= 32 and image.size[1] >= 32:
    #                         embedding = self.process_image(image)
    #                         documents.append({
    #                             'content': embedding,
    #                             'type': 'image',
    #                             'source': f"{docx_path}#image{img_idx}"
    #                         })
    #                 except Exception as e:
    #                     self.logger.error(f"Error extracting image {img_idx} from {docx_path}: {str(e)}")
    #                     continue
            
    #         self.clear_memory()
        
    #     except Exception as e:
    #         self.logger.error(f"Error processing DOCX {docx_path}: {str(e)}")
        
    #     return documents
    
    def process_project(self, project_name: str, project_dir: str,sow_data):
        """Process all files in a project directory with memory management"""
        print(f"\n==== Processing project: {project_name} ====\nDirectory: {project_dir}")
        
        # Create collection
        # try:
        #     collection = self.create_project_collection(project_name)
        #     print(f"ChromaDB collection created/accessed successfully")
        # except Exception as e:
        #     print(f"ERROR creating ChromaDB collection: {str(e)}")
        #     raise
            
        documents = {}
        
        # Get list of files
        files = []
        for root, _, filenames in os.walk(project_dir):
            for file in filenames:
                file_path = os.path.join(root, file)
                files.append(file_path)
                
        print(f"Found {len(files)} files to process: {[os.path.basename(f) for f in files]}")
        
        for file_path in files:
            try:
                documents[file_path]=self.parse_file(file_path,sow_data)
                self.logger.info(f"Processed image file")
            except Exception as e:
                self.logger.error(f"Error processing file {file_path}: {str(e)}")
            finally:
                self.clear_memory()
        
        print(documents)
        self.clear_memory()
        import json
        with open("documents.json", "w") as f:
            json.dump(documents, f, indent=4)
        
        if len(documents) >= 50:
            print(f"Processing batch of {len(documents)} documents")
            # self.process_documents_batch(documents, collection)
            documents = []
            self.clear_memory()
    
        if documents:
            print(f"Processing final batch of {len(documents)} documents")
            # self.process_documents_batch(documents, collection)
            self.clear_memory()
            
        print(f"\n==== Completed processing project: {project_name} ====\n")

    def process_documents_batch(self, documents: List[Dict], collection=None):
        """Process a batch of documents and add to collection"""
        if documents:
            print(f"Creating embeddings for {len(documents)} documents")
            self.logger.info(f"Creating embeddings for {len(documents)} documents")
            
            embed_data = self.create_embeddings(documents)
            print(f"Successfully created {len(embed_data['embeddings'])} embeddings")
            if embed_data['embeddings']:
                try:
                    first_embedding_dim = len(embed_data['embeddings'][0])
                    self.logger.info(f"Embedding dimension: {first_embedding_dim}")
                    if first_embedding_dim != 384:
                        raise ValueError(f"Expected embedding dimension 384, got {first_embedding_dim}")
                    
                    batch_size = 20
                    for i in range(0, len(embed_data['embeddings']), batch_size):
                        batch_end = min(i + batch_size, len(embed_data['embeddings']))
                        
                        batch_documents = []
                        for j in range(i, batch_end):
                            if embed_data['metadatas'][j]['type'] == 'text':
                                orig_doc = next(
                                (doc for doc in documents if doc['source'] == embed_data['metadatas'][j]['source']),
                                None
                                )
                                batch_documents.append(orig_doc['content'] if orig_doc else None)
                            else:
                                batch_documents.append(None)
                        
                        try:
                            print(f"Adding batch of {batch_end - i} embeddings to ChromaDB collection")
                            collection.add(
                                embeddings=embed_data['embeddings'][i:batch_end],
                                metadatas=embed_data['metadatas'][i:batch_end],
                                    documents=batch_documents,
                                    ids=embed_data['ids'][i:batch_end]
                                )
                            print(f"Successfully added embeddings to ChromaDB")
                            self.logger.info(f"Added batch of {batch_end - i} embeddings to collection")
                            self.clear_memory()
                        except Exception as add_error:
                            print(f"ERROR adding embeddings to ChromaDB: {str(add_error)}")
                            raise
                
                except Exception as e:
                    self.logger.error(f"Error adding embeddings to collection: {str(e)}")
                    if embed_data['embeddings']:
                        self.logger.error(f"First embedding shape: {len(embed_data['embeddings'][0])}")

    def get_document_content(self, source: str) -> str:
        """Retrieve original document content from source path"""
        try:
            if '#page' in source:  # PDF document
                pdf_path, page_info = source.split('#page')
                page_num = int(page_info) - 1
                with fitz.open(pdf_path) as doc:
                    return doc[page_num].get_text()
            elif '#slide' in source:  # PPT document
                ppt_path, slide_info = source.split('#slide')
                slide_num = int(slide_info) - 1
                presentation = Presentation(ppt_path)
                text = ""
                for shape in presentation.slides[slide_num].shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                return text
            elif '#para' in source or '#table' in source:  # DOCX document
                docx_path = source.split('#')[0]
                doc = docx.Document(docx_path)
                
                if '#para' in source:
                    # Get paragraph chunk
                    para_idx = int(source.split('#para')[1])
                    chunk_size = 3
                    start_idx = para_idx * chunk_size
                    end_idx = min(start_idx + chunk_size, len(doc.paragraphs))
                    paragraphs = [p.text for p in doc.paragraphs[start_idx:end_idx] if p.text.strip()]
                    return "\n".join(paragraphs)
                    
                elif '#table' in source:
                    # Get table
                    table_idx = int(source.split('#table')[1])
                    if table_idx < len(doc.tables):
                        table = doc.tables[table_idx]
                        table_text = []
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                row_text.append(cell.text)
                            table_text.append(" | ".join(row_text))
                        return "\n".join(table_text)
                    
                return "Content not available"
            return "Content not available"
        except Exception as e:
            self.logger.error(f"Error retrieving document content: {str(e)}")
            return "Error retrieving content"
    def parse_file(self, file_path, sow_data):
        print(f"\nProcessing file: {file_path}")
        try:
            self.logger.info(f"Processing file: {file_path}")
            self.logger.info(f"Current RAM usage: {get_memory_usage():.2f} MB")
            
            # If file is not already a PDF, convert it
            if file_path.lower().endswith(('.docx', '.doc', '.ppt', '.pptx')):
                print(f"Converting {file_path} to PDF")
                pdf_path = convert_to_pdf(file_path)
                if pdf_path:
                    print(f"Conversion successful: {pdf_path}")
            else:
                pdf_path = file_path if file_path.lower().endswith('.pdf') else None
            
            # Process based on file type
            if pdf_path and pdf_path.lower().endswith('.pdf'):
                try:
                    # Extract requirements from SOW data
                    requirements = sow_data.get('requirements', [])
                    
                    # Extract content with Gemini
                    extracted_info = extract_text_with_gemini(pdf_path, self.gemini_api_key, sow_data)
                    
                    # NEW: Match requirements to document content
                    requirement_matches = match_requirements_to_document(requirements, extracted_info)
                    
                    # Add matches to the extracted info
                    result = {
                        'extracted_content': extracted_info,
                        'requirement_matches': requirement_matches,
                        'source_file': file_path
                    }
                    
                    return result
                except Exception as e:
                    self.logger.error(f"Error processing with Gemini: {str(e)}")      
                
        except Exception as e:
            self.logger.error(f"Error processing file {file_path}: {str(e)}")
        finally:
            self.clear_memory()

    def generate_document_clusters(self, project_name: str, eps: float = None, min_samples: int = None) -> Dict[str, Any]:
        """Generate and store document clusters for a project
        
        Args:
            project_name: Name of the project collection
            eps: DBSCAN epsilon parameter (auto-calculated if None)
            min_samples: DBSCAN min_samples parameter (auto-calculated if None)
            
        Returns:
            Dictionary with cluster information
        """
        print(f"\n==== Generating document clusters for '{project_name}' ====")
        
        try:
            # Get the collection
            collection = self.chroma_client.get_collection(project_name)
            print(f"Successfully connected to collection '{project_name}'")
            
            # Retrieve document embeddings
            results = collection.get(include=["embeddings", "metadatas", "documents"])
            embeddings = results["embeddings"]
            metadatas = results["metadatas"]
            documents = results["documents"]
            
            print(f"Retrieved {len(embeddings)} document embeddings")
            import numpy as np
            # Convert to numpy array
            embeddings_array = np.array(embeddings, dtype=np.float32)
            
            # Check if the array is empty
            if embeddings_array.size == 0:
                return {"error": "No embeddings found", "clusters": {}}
            
            # Auto-calculate DBSCAN parameters if not provided
            if eps is None:
                # Calculate pairwise distances and use the 90th percentile as eps
                from sklearn.metrics import pairwise_distances
                import numpy as np
                
                # Use a sample if there are many documents
                max_sample = 1000
                if len(embeddings_array) > max_sample:
                    indices = np.random.choice(len(embeddings_array), max_sample, replace=False)
                    sample = embeddings_array[indices]
                else:
                    sample = embeddings_array
                
                # Calculate pairwise distances
                distances = pairwise_distances(sample, metric='cosine')
                # Flatten the upper triangle
                dist_list = distances[np.triu_indices_from(distances, k=1)]
                # Use the 10th percentile as eps (smaller for tighter clusters)
                eps = np.percentile(dist_list, 25)
                print(f"Auto-calculated eps value: {eps}")
            
            if min_samples is None:
                # Rule of thumb: min_samples = 2 * dimensions
                min_samples = min(5, int(embeddings_array.shape[1] * 0.05))
                print(f"Auto-calculated min_samples value: {min_samples}")
            
            # Perform DBSCAN clustering
            print(f"Running DBSCAN with eps={eps}, min_samples={min_samples}...")
            from sklearn.cluster import DBSCAN
            
            # Normalize embeddings
            norms = np.linalg.norm(embeddings_array, axis=1, keepdims=True)
            norms[norms == 0] = 1  # Avoid division by zero
            normalized_embeddings = embeddings_array / norms
            
            dbscan = DBSCAN(eps=eps, min_samples=min_samples, metric='cosine')
            clusters = dbscan.fit_predict(normalized_embeddings)
            
            # Count clusters (excluding noise points labeled as -1)
            unique_clusters = set(clusters)
            n_clusters = len(unique_clusters) - (1 if -1 in unique_clusters else 0)
            n_noise = list(clusters).count(-1)
            
            print(f"DBSCAN found {n_clusters} clusters, {n_noise} noise points")
            
            # Organize documents by cluster
            cluster_data = {}
            for cluster_id in unique_clusters:
                if cluster_id == -1:
                    cluster_name = "Outliers"
                else:
                    cluster_name = f"Cluster_{cluster_id}"
                
                # Get documents in this cluster
                cluster_indices = np.where(clusters == cluster_id)[0]
                
                # Extract cluster documents and metadata
                cluster_docs = []
                for idx in cluster_indices:
                    if idx < len(documents) and documents[idx]:
                        doc_data = {
                            'content': documents[idx],
                            'source': metadatas[idx].get('source', ''),
                            'type': metadatas[idx].get('type', '')
                        }
                        cluster_docs.append(doc_data)
                
                # Generate cluster summary
                summary = self._summarize_cluster(cluster_docs)
                
                cluster_data[cluster_name] = {
                    'document_count': len(cluster_docs),
                    'documents': cluster_docs,
                    'summary': summary
                }
            
            # Update document metadata with cluster information
            print("Updating document metadata with cluster information...")
            # Batched updates to avoid memory issues
            batch_size = 20
            for i in range(0, len(clusters), batch_size):
                batch_end = min(i + batch_size, len(clusters))
                batch_ids = [f"doc_{j}" for j in range(i, batch_end)]
                
                # Get existing metadata
                batch_metadatas = []
                for j in range(i, batch_end):
                    if j < len(metadatas):
                        metadata = metadatas[j].copy()
                        cluster_id = int(clusters[j])
                        if cluster_id == -1:
                            metadata['cluster'] = "Outliers"
                        else:
                            metadata['cluster'] = f"Cluster_{cluster_id}"
                        batch_metadatas.append(metadata)
                
                # Update metadata in the collection
                if batch_metadatas:
                    collection.update(
                        ids=batch_ids,
                        metadatas=batch_metadatas
                    )
            
            return {
                "status": "success",
                "clusters": cluster_data,
                "parameters": {
                    "eps": eps,
                    "min_samples": min_samples
                }
            }
        
        except Exception as e:
            print(f"Error generating document clusters: {str(e)}")
            import traceback
            traceback.print_exc()
            return {"error": str(e), "clusters": {}}
        
    def _summarize_cluster(self, cluster_docs: List[Dict[str, Any]]) -> str:
        """Generate a summary for a document cluster"""
        try:
            # Extract text content
            text_content = []
            for doc in cluster_docs:
                if doc.get('type') == 'text' and doc.get('content'):
                    text_content.append(doc.get('content'))
            
            if not text_content:
                return "No text content available for summarization"
            
            # Combine text (limit to avoid token limits)
            combined_text = "\n\n".join(text_content)
            if len(combined_text) > 10000:
                combined_text = combined_text[:10000] + "... (truncated)"
            
            # Use Gemini to summarize if available
            try:
                import google.generativeai as genai
                model = genai.GenerativeModel('gemini-2.0-flash')
                
                prompt = f"""
                Summarize the main themes and topics in these related documents.
                Be concise and focus on the key information.
                
                Documents:
                {combined_text}
                
                Summary:
                """
                
                response = model.generate_content(prompt)
                return response.text.strip()
            except Exception as ge:
                print(f"Error using Gemini for summarization: {str(ge)}")
                # Fall back to simple summarization
                words = combined_text.split()
                return " ".join(words[:100]) + "..." if len(words) > 100 else combined_text
        
        except Exception as e:
            print(f"Error summarizing cluster: {str(e)}")
            return "Error generating summary"

    def list_projects(self) -> List[str]:
        try:
            conn = sqlite3.connect('discovery.db')
            df = pd.read_sql("Select * from projects", conn)
            conn.close()
            return df['name'].tolist()

        except:
            pass

    def get_project_stats(self, project_name: str) -> Dict:
        """Get statistics about a specific project's collection"""
        try:
            collection = self.chroma_client.get_collection(name=project_name)
            count = collection.count()
            
            results = collection.get(
                include=["metadatas"]
            )
            
            text_count = sum(1 for meta in results['metadatas'] if meta['type'] == 'text')
            image_count = sum(1 for meta in results['metadatas'] if meta['type'] == 'image')
            
            return {
                'total_documents': count,
                'text_documents': text_count,
                'image_documents': image_count
            }
        except Exception as e:
            self.logger.error(f"Error getting stats for project {project_name}: {str(e)}")
            return {}
        
    def summarize_existing_chroma(self, collection_name=None, eps=0.5, min_samples=5, gemini_model="gemini-2.0-flash"):
        print(f"\n==========================================\nSTARTING SUMMARIZATION FOR: {collection_name}\n==========================================")
        """Summarize an existing Chroma DB collection using DBSCAN clustering and Gemini with enhanced logging.
        
        Args:
            collection_name: Name of the Chroma collection to summarize
            eps: DBSCAN parameter - maximum distance between points in the same neighborhood
            min_samples: DBSCAN parameter - minimum number of points to form a dense region
            gemini_model: Gemini model to use for summarization
        
        Returns:
            A final summary of the collection
        """
        print(f"\n==== Summarizing collection '{collection_name}' with DBSCAN (eps={eps}, min_samples={min_samples}) ====\nTIMESTAMP: {time.strftime('%Y-%m-%d %H:%M:%S')}")

        # 1. Get the collection using the existing client
        try:
            collection = self.chroma_client.get_collection(collection_name)
            print(f"Successfully connected to collection '{collection_name}'")
        except Exception as e:
            print(f"Error getting collection: {e}")
            return f"Error: Could not access collection '{collection_name}'"
        
        # 2. Retrieve all documents
        try:
            results = collection.get(include=["embeddings", "documents"])
            documents = results["documents"]
            embeddings = results["embeddings"]
            print(f"Retrieved {len(documents)} documents with embeddings from collection")
        except Exception as e:
            print(f"Error retrieving documents: {e}")
            return f"Error: Failed to retrieve documents from collection: {e}"
        
        # 3. Validate and prepare embeddings for clustering
        try:
            # Convert to numpy array and check for NaN values
            embeddings_array = np.array(embeddings, dtype=np.float32)
            
            # Check if the array is empty
            if embeddings_array.size == 0:
                return "Error: No embeddings found in the collection"
                
            # Check for NaN values
            if np.isnan(embeddings_array).any():
                print("Warning: NaN values found in embeddings. Removing documents with NaN embeddings.")
                # Create a mask for rows without NaN values
                valid_mask = ~np.isnan(embeddings_array).any(axis=1)
                embeddings_array = embeddings_array[valid_mask]
                documents = [doc for i, doc in enumerate(documents) if valid_mask[i]]
                print(f"After removing NaN values: {len(documents)} documents remaining")
                
                if len(documents) == 0:
                    return "Error: All embeddings contain NaN values"
            
            # Ensure we have a 2D array
            if len(embeddings_array.shape) == 1:
                print(f"Reshaping 1D array of shape {embeddings_array.shape} to 2D")
                embeddings_array = embeddings_array.reshape(1, -1)
            
            print(f"Embeddings array shape: {embeddings_array.shape}")
        except Exception as e:
            print(f"Error processing embeddings: {e}")
            return f"Error: Failed to process embeddings: {e}"
        
        # 4. Perform DBSCAN clustering
        try:
            print(f"Clustering embeddings using DBSCAN with eps={eps}, min_samples={min_samples}...")
            print(f"Embeddings array shape before DBSCAN: {embeddings_array.shape}")
            print(f"Embeddings array data type: {embeddings_array.dtype}")
            print(f"First embedding vector sample (first 5 values): {embeddings_array[0][:5]}")
            
            # Check for infinite values
            if np.isinf(embeddings_array).any():
                print("WARNING: Infinite values found in embeddings. Replacing with zeros.")
                embeddings_array = np.nan_to_num(embeddings_array, nan=0.0, posinf=0.0, neginf=0.0)
            
            # Log memory usage
            print(f"Memory usage of embeddings array: {embeddings_array.nbytes / (1024 * 1024):.2f} MB")
            
            from sklearn.cluster import DBSCAN
            try:
                # Normalize embeddings to ensure uniform distance measures
                print("Normalizing embeddings for consistent distance calculations...")
                norms = np.linalg.norm(embeddings_array, axis=1, keepdims=True)
                # Avoid division by zero
                norms[norms == 0] = 1
                normalized_embeddings = embeddings_array / norms
                print("Normalization complete.")
                
                # Initialize and fit DBSCAN
                print(f"Running DBSCAN with eps={eps}, min_samples={min_samples}...")
                dbscan = DBSCAN(eps=eps, min_samples=min_samples)
                clusters = dbscan.fit_predict(normalized_embeddings)
                print("DBSCAN completed successfully.")
            except Exception as dbscan_err:
                print(f"Error during DBSCAN execution: {dbscan_err}")
                print("Trying DBSCAN with original embeddings without normalization...")
                dbscan = DBSCAN(eps=eps, min_samples=min_samples)
                clusters = dbscan.fit_predict(embeddings_array)
                print("DBSCAN completed with original embeddings.")
            
            # Count number of clusters (excluding noise points labeled as -1)
            unique_clusters = set(clusters)
            n_clusters = len(unique_clusters) - (1 if -1 in unique_clusters else 0)
            n_noise = list(clusters).count(-1)
            
            print(f"DBSCAN clustering completed: {n_clusters} clusters found, {n_noise} noise points")
            
            # Log cluster distribution
            if n_clusters > 0:
                print("Cluster distribution:")
                for cluster_id in unique_clusters:
                    if cluster_id != -1:  # Skip noise points
                        count = list(clusters).count(cluster_id)
                        print(f"  Cluster {cluster_id}: {count} points ({count/len(clusters)*100:.1f}%)")
            
            # If no clusters were found, try with different parameters
            if n_clusters == 0:
                print("No clusters found. Trying with different parameters...")
                new_eps = eps * 1.5  # Increase eps by 50%
                new_min_samples = max(2, min_samples - 1)  # Decrease min_samples but keep at least 2
                print(f"New parameters: eps={new_eps}, min_samples={new_min_samples}")
                
                dbscan = DBSCAN(eps=new_eps, min_samples=new_min_samples)
                clusters = dbscan.fit_predict(embeddings_array)
                unique_clusters = set(clusters)
                n_clusters = len(unique_clusters) - (1 if -1 in unique_clusters else 0)
                n_noise = list(clusters).count(-1)
                print(f"Retry DBSCAN: {n_clusters} clusters found, {n_noise} noise points with eps={new_eps}, min_samples={new_min_samples}")
                eps = new_eps  # Update eps for later use
                min_samples = new_min_samples
            
            # If still no clusters, fall back to a very inclusive approach
            if n_clusters == 0:
                print("Still no clusters found. Using a very inclusive approach...")
                eps = 0.9  # Very large eps
                min_samples = 2  # Minimum possible value
                print(f"Last attempt parameters: eps={eps}, min_samples={min_samples}")
                
                dbscan = DBSCAN(eps=eps, min_samples=min_samples)
                clusters = dbscan.fit_predict(embeddings_array)
                unique_clusters = set(clusters)
                n_clusters = len(unique_clusters) - (1 if -1 in unique_clusters else 0)
                n_noise = list(clusters).count(-1)
                print(f"Final DBSCAN attempt: {n_clusters} clusters found, {n_noise} noise points")
            
            # If still no clusters, create a fallback approach
            if n_clusters == 0:
                print("Could not identify any clusters. Creating a fallback approach.")
                # Treat each document as its own cluster for the top 10 documents
                clusters = np.ones(len(documents), dtype=int) * -1  # All noise initially
                max_docs = min(10, len(documents))
                for i in range(max_docs):
                    clusters[i] = i  # Assign each of the first max_docs to its own cluster
                unique_clusters = set(clusters)
                n_clusters = len(unique_clusters) - (1 if -1 in unique_clusters else 0)
                n_noise = len(documents) - max_docs
                print(f"Using fallback approach: Created {n_clusters} manual clusters with {n_noise} noise points")
        except Exception as e:
            print(f"Error during clustering: {e}")
            import traceback
            traceback.print_exc()
            return f"Error: Clustering failed: {e}"
        
        # 5. Find representative documents for each cluster
        try:
            print("Finding representative documents for each cluster...")
            representative_docs = []
            cluster_indices = []
            
            # Get all unique cluster IDs (excluding -1 which is noise)
            cluster_ids = [c for c in set(clusters) if c != -1]
            
            # Process each cluster
            for cluster_id in cluster_ids:
                # Get indices of documents in this cluster
                cluster_docs_indices = np.where(clusters == cluster_id)[0]
                
                if len(cluster_docs_indices) > 0:
                    # For text clustering, select a few representative documents from each cluster
                    max_representatives = min(3, len(cluster_docs_indices))
                    # Safely sample with replacement if needed
                    replace = len(cluster_docs_indices) < max_representatives
                    sampled_indices = np.random.choice(cluster_docs_indices, max_representatives, replace=replace)
                    
                    for idx in sampled_indices:
                        if idx < len(documents) and documents[idx] and documents[idx].strip():  # Ensure document has content
                            representative_docs.append(documents[idx])
                            cluster_indices.append(cluster_id)
                    
                    # Print cluster statistics
                    print(f"Cluster {cluster_id}: {len(cluster_docs_indices)} documents, selected {len(sampled_indices)} representatives")
            
            # Process noise points separately if we have them
            noise_indices = np.where(clusters == -1)[0]
            if len(noise_indices) > 0:
                # Include some noise points if we have too few representative documents or if specifically requested
                if len(representative_docs) < 5:
                    max_noise_docs = min(5 - len(representative_docs), len(noise_indices))
                    
                    if max_noise_docs > 0:
                        replace = len(noise_indices) < max_noise_docs
                        sampled_noise = np.random.choice(noise_indices, max_noise_docs, replace=replace)
                        
                        for idx in sampled_noise:
                            if idx < len(documents) and documents[idx] and documents[idx].strip():  # Ensure document has content
                                representative_docs.append(documents[idx])
                                cluster_indices.append(-1)  # Indicating this is from noise
                        
                        print(f"Added {max_noise_docs} noise documents to ensure sufficient representation")
            
            print(f"Selected {len(representative_docs)} representative documents from {n_clusters} clusters")
            
            # Fallback if no representative docs were found
            if len(representative_docs) == 0:
                print("No representative documents found. Using fallback approach.")
                # Use the first few documents as representatives
                max_docs = min(5, len(documents))
                for i in range(max_docs):
                    if i < len(documents) and documents[i] and documents[i].strip():
                        representative_docs.append(documents[i])
                        cluster_indices.append(0)  # Assign to a dummy cluster
                
                print(f"Using {len(representative_docs)} documents as fallback representatives")
        except Exception as e:
            print(f"Error finding representative documents: {e}")
            import traceback
            traceback.print_exc()
            return f"Error: Failed to find representative documents: {e}"
        
        # 6. Summarize representative documents using Gemini
        try:
            print("Summarizing representative documents with Gemini...")
            model = genai.GenerativeModel(gemini_model)
            
            summaries = []
            for i, doc in enumerate(representative_docs):
                # Skip empty documents
                if not doc or doc.strip() == "":
                    print(f"Skipping empty document from cluster {cluster_indices[i]}")
                    continue
                    
                # Truncate very long documents
                max_length = 8000
                if len(doc) > max_length:
                    print(f"Truncating document from {len(doc)} to {max_length} characters")
                    doc = doc[:max_length] + "... [truncated]"
                
                prompt = f"""
                Summarize the following document in a concise manner, focusing on the main points and key details:
                
                {doc}
                
                Summary:
                """
                
                try:
                    response = model.generate_content(prompt)
                    summary = response.text
                    cluster_label = f"Cluster {cluster_indices[i]}" if cluster_indices[i] != -1 else "Noise points"
                    summaries.append(f"{cluster_label}:\n{summary}")
                except Exception as e:
                    print(f"Error summarizing document {i}: {e}")
                    summaries.append(f"Error generating summary for document in cluster {cluster_indices[i]}.")
        except Exception as e:
            print(f"Error during document summarization: {e}")
            return f"Error: Document summarization failed: {e}"
        
        # 7. Create a final combined summary with Gemini
        try:
            print(f"\n==== STEP 7: Creating final combined summary at {time.strftime('%Y-%m-%d %H:%M:%S')} ====")
            print("Creating final combined summary...")
            combined_summaries = "\n\n".join(summaries)
            
            final_summary_prompt = f"""
            Below are summaries of different clusters of documents from a vector database.
            Create a cohesive, comprehensive summary that captures the key information across all these summaries.
            Organize the information logically and remove any redundancies.
            
            CLUSTER SUMMARIES:
            {combined_summaries}
            
            FINAL COMPREHENSIVE SUMMARY:
            """
            
            try:
                final_response = model.generate_content(final_summary_prompt)
                final_summary = final_response.text
            except Exception as e:
                print(f"Error creating final summary: {e}")
                final_summary = "Error generating final summary. Individual cluster summaries:\n\n" + combined_summaries
        except Exception as e:
            print(f"Error during final summarization: {e}")
            return f"Error: Final summarization failed: {e}"
        
        print(f"\n==== COMPLETED SUMMARIZATION FOR: {collection_name} at {time.strftime('%Y-%m-%d %H:%M:%S')} ====\n")
        return final_summary


def convert_ppt_to_pdf(input_file, output_file=None):
    """Convert a PowerPoint file to PDF"""
    convert(input_file, output_file)

def convert_to_pdf(file_path):
    """Convert a document (PPTX or DOCX) to PDF"""
    extension = os.path.splitext(file_path)[1].lower()
    output_file = os.path.splitext(file_path)[0] + ".pdf"
    
    if extension in ['.docx', '.doc']:
        docx_convert(file_path, output_file)
    elif extension in ['.pptx', '.ppt']:
        convert_ppt_to_pdf(file_path, output_file)
    else:
        return None  # No conversion needed or not supported
    
    return output_file

def match_requirements_to_document( requirements, document_content):
    """Match requirements to document content to find supporting sections"""
    requirement_matches = {}
    
    try:
        # Convert document content to string if needed
        if isinstance(document_content, dict):
            document_text = json.dumps(document_content)
        else:
            document_text = str(document_content)
        
        for req in requirements:
            req_id = req.get('id', '')
            req_text = req.get('text', '')
            
            # Skip empty requirements
            if not req_text.strip():
                continue
            
            # Generate keywords and phrases from the requirement
            keywords = extract_keywords(req_text)
            
            # Find matches in document text
            matches = []
            for keyword in keywords:
                # Skip very short keywords (likely to cause false positives)
                if len(keyword) < 4:
                    continue
                    
                # Find keyword occurrences
                pattern = re.compile(r'(.{0,100}' + re.escape(keyword) + r'.{0,100})', re.IGNORECASE)
                for match in pattern.finditer(document_text):
                    # Extract the context around the keyword
                    context = match.group(1).strip()
                    if context:
                        matches.append({
                            'keyword': keyword,
                            'context': context
                        })
            
            # Store matches for this requirement
            if matches:
                requirement_matches[req_id] = matches
        
        return requirement_matches
    
    except Exception as e:
        print(f"Error matching requirements to document: {str(e)}")
        return {}
    
def extract_keywords(text):
    """Extract keywords and phrases from requirement text"""
    keywords = []
    
    # Remove common stop words
    stop_words = ['the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 'with', 'by', 'shall', 'will', 'should', 'must']
    
    # Simple keyword extraction
    words = text.lower().split()
    filtered_words = [word.strip(',.()[]{}:;\'\"') for word in words if word.lower() not in stop_words and len(word) > 3]
    keywords.extend(filtered_words)
    
    # Extract phrases (2-3 word combinations)
    for i in range(len(words) - 1):
        phrase = words[i] + ' ' + words[i+1]
        keywords.append(phrase.strip(',.()[]{}:;\'\"'))
    
    for i in range(len(words) - 2):
        phrase = words[i] + ' ' + words[i+1] + ' ' + words[i+2]
        keywords.append(phrase.strip(',.()[]{}:;\'\"'))
    
    # Remove duplicates
    keywords = list(set(keywords))
    
    return keywords

def extract_text_with_gemini(pdf_path, gemini_api_key, sow_data):
    """Use Gemini to extract structured information from a PDF"""
    genai.configure(api_key=gemini_api_key)
    
    # Use Gemini Pro Vision model for processing PDFs
    model = genai.GenerativeModel('gemini-2.0-flash')
    
    # Read PDF as bytes for the API
    with open(pdf_path, 'rb') as f:
        pdf_bytes = f.read()
    
    # Extract requirements information for context
    requirements_text = ""
    if sow_data and 'requirements' in sow_data:
        for req in sow_data['requirements'][:10]:  # Limit to avoid token limits
            requirements_text += f"- {req.get('id', '')}: {req.get('text', '')}\n"
    
    # Create a prompt that extracts key information and tries to match requirements
    prompt = f"""
    Extract all the text and the understanding from this document.
    For diagrams properly extract the text and the understanding from it.
    
    Additionally, I'm providing a list of key requirements from the SOW. 
    In your analysis, please identify any content in this document that relates to these requirements:
    
    {requirements_text}
    
    Format the output as structured JSON with:
    1. Clear sections of the document content
    2. Any key information related to the requirements
    3. Important technical details and specifications
    """
    
    # Create the image part from PDF bytes (first page)
    # Note: For multi-page PDFs, you may need to split it first
    response = model.generate_content([
        prompt,
        {'mime_type': 'application/pdf', 'data': pdf_bytes}
    ])
    
    # Process and return the extraction result
    return response.text